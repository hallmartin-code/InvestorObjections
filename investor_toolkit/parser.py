"""Deck ingestion: extract plain text from PDF or PPTX pitch decks."""

from __future__ import annotations

import re
from pathlib import Path

SLIDE_BREAK = "\n---\n"
SUPPORTED_EXTENSIONS = (".pdf", ".pptx", ".docx")


def extract_text(filepath: str) -> str:
    """Extract the full text of a pitch deck, preserving slide/page order.

    Slides and pages are separated by a ``\\n---\\n`` marker so downstream prompts can
    still see the deck's structure.

    Raises:
        FileNotFoundError: the deck does not exist.
        ValueError: the file extension is not ``.pdf`` or ``.pptx``.
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"Deck not found: {path}")
    if not path.is_file():
        raise ValueError(f"Not a file: {path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        pages = _extract_pdf(path)
    elif suffix == ".pptx":
        pages = _extract_pptx(path)
    elif suffix == ".docx":
        pages = _extract_docx(path)
    else:
        raise ValueError(
            f"Unsupported file type '{suffix or path.name}'. "
            f"Supported types: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    cleaned = [_clean(p) for p in pages]
    return SLIDE_BREAK.join(p for p in cleaned if p)


def _extract_pdf(path: Path) -> list[str]:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            pages.append(page.extract_text() or "")
    return pages


def _extract_pptx(path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[str] = []
    for slide in prs.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            chunks.extend(_shape_text(shape))
        slides.append("\n".join(c for c in chunks if c.strip()))
    return slides


def _extract_docx(path: Path) -> list[str]:
    """Walk the document body in order, splitting on explicit page breaks."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    def has_page_break(element) -> bool:
        if any(node.get(qn("w:type")) == "page" for node in element.iter(qn("w:br"))):
            return True
        # Word writes this where it laid out a page; useful for decks authored as documents.
        return next(element.iter(qn("w:lastRenderedPageBreak")), None) is not None

    document = Document(str(path))
    body = document.element.body
    pages: list[list[str]] = [[]]

    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            para = Paragraph(child, document)
            text = para.text.strip()
            if text:
                pages[-1].append(text)
            # A page break starts a new chunk, mirroring PDF/PPTX pagination.
            if has_page_break(child):
                pages.append([])
        elif tag == "tbl":
            for row in Table(child, document).rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    pages[-1].append(line)

    return ["\n".join(chunk) for chunk in pages if chunk]


def _shape_text(shape) -> list[str]:
    """Collect text from a shape, recursing into groups and reading table cells."""
    out: list[str] = []

    # Group shapes (MSO_SHAPE_TYPE.GROUP == 6) hold their own child shapes.
    if getattr(shape, "shape_type", None) == 6:
        for child in shape.shapes:
            out.extend(_shape_text(child))
        return out

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                out.append(line)
        return out

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                out.append(text)
        return out

    text = getattr(shape, "text", "")
    if isinstance(text, str) and text.strip():
        out.append(text.strip())
    return out


def _clean(text: str) -> str:
    """Strip excessive whitespace while keeping paragraph structure."""
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\xa0", " ")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
